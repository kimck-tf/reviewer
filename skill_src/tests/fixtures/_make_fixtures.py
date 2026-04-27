"""
Fixture PPTX 생성 스크립트.
한 번 실행하여 fixtures/*.pptx를 생성. 결과는 git에 commit하여 재실행 불필요.

sample_with_embedded.pptx는 PowerPoint UI에서 수동 생성해야 함
(임베디드 OLE 객체는 python-pptx로 생성 불가).
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


HERE = Path(__file__).parent


def make_text_only() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Slide 1: 표지
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    slide.shapes.title.text = "프론트 서스펜션 강도 해석"
    slide.placeholders[1].text = "작성자: 홍길동\n2026-04-27"

    # Slide 2: 목차
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "목차"
    slide.placeholders[1].text = (
        "1. 해석 목적\n2. 모델 정보\n3. 경계조건\n4. 결과\n5. 결론"
    )

    # Slide 3~5: 본문
    for i, t in enumerate(["해석 목적", "모델 정보", "결과"]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = t
        slide.placeholders[1].text = f"{t}에 대한 설명입니다.\n참고 사항을 적습니다."
        # 발표자 노트
        slide.notes_slide.notes_text_frame.text = f"슬라이드 {i+3} 발표 노트"

    prs.save(HERE / "sample_text_only.pptx")
    print(f"Created: {HERE / 'sample_text_only.pptx'}")


def make_with_table() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide.shapes.title.text = "결과 표"

    # 표 추가
    rows, cols = 4, 3
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["항목", "값", "단위"]
    data = [
        ["최대 응력", "240", "MPa"],
        ["허용 응력", "350", "MPa"],
        ["안전계수", "1.46", "-"],
    ]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, row_data in enumerate(data, start=1):
        for c, v in enumerate(row_data):
            table.cell(r, c).text = v

    prs.save(HERE / "sample_with_table.pptx")
    print(f"Created: {HERE / 'sample_with_table.pptx'}")


def make_with_image() -> None:
    """단순 도형(차트 대신)을 추가한 슬라이드. 실제 이미지 파일은 PIL로 즉석 생성."""
    from pptx.util import Emu
    from PIL import Image
    import io

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "그래프"

    # 더미 이미지 생성
    img = Image.new("RGB", (400, 300), color=(200, 220, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(2), Inches(2), Inches(5), Inches(3.75))

    prs.save(HERE / "sample_with_image.pptx")
    print(f"Created: {HERE / 'sample_with_image.pptx'}")


if __name__ == "__main__":
    make_text_only()
    make_with_table()
    make_with_image()
    print("\nNOTE: sample_with_embedded.pptx는 PowerPoint에서 수동 생성:")
    print("  1. 새 PPT 만들기")
    print("  2. 슬라이드에 '삽입 > 개체 > Microsoft Excel Worksheet' 추가")
    print("  3. fixtures/sample_with_embedded.pptx로 저장")
