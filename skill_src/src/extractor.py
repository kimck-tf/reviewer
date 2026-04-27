from __future__ import annotations
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract(pptx_path: Path) -> dict[str, Any]:
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width or 0
    slide_h = prs.slide_height or 0

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        shapes = _extract_shapes(slide, slide_index=idx, slide_w=slide_w, slide_h=slide_h)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text or ""
        slides.append({
            "index": idx,
            "title": _get_slide_title(slide),
            "layout": slide.slide_layout.name if slide.slide_layout else None,
            "shapes": shapes,
            "notes": notes,
            "image_path": None,
            "thumbnail_path": None,
            "embedded_image_paths": [],
            "has_embedded": False,
        })

    cp = prs.core_properties
    metadata = {
        "file_path": str(pptx_path),
        "title": cp.title or "",
        "author": cp.author or "",
        "created": cp.created.isoformat() if cp.created else "",
        "modified": cp.modified.isoformat() if cp.modified else "",
        "slide_count": len(slides),
    }
    return {"metadata": metadata, "slides": slides}


def _get_slide_title(slide) -> str:
    if slide.shapes.title is not None:
        return slide.shapes.title.text or ""
    return ""


def _extract_shapes(slide, slide_index: int, slide_w: int, slide_h: int) -> list[dict[str, Any]]:
    out = []
    for sh_idx, shape in enumerate(slide.shapes, start=1):
        shape_type = _classify_shape_type(shape)
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text or ""

        table_data = None
        if shape_type == "Table" and shape.has_table:
            tbl = shape.table
            cells = [[cell.text for cell in row.cells] for row in tbl.rows]
            table_data = {
                "rows": len(tbl.rows),
                "cols": len(tbl.columns),
                "cells": cells,
            }

        left = int(shape.left) if shape.left is not None else 0
        top = int(shape.top) if shape.top is not None else 0
        width = int(shape.width) if shape.width is not None else 0
        height = int(shape.height) if shape.height is not None else 0

        position_pct: dict[str, float] = {}
        if slide_w and slide_h:
            position_pct = {
                "left": left / slide_w,
                "top": top / slide_h,
                "width": width / slide_w,
                "height": height / slide_h,
            }

        out.append({
            "shape_id": f"s{slide_index}_sh{sh_idx}",
            "type": shape_type,
            "position_emu": {"left": left, "top": top, "width": width, "height": height},
            "position_pct": position_pct,
            "z_order": sh_idx,
            "text": text,
            "table": table_data,
            "image_ref": None,
            "embedded_progid": None,
        })
    return out


def _classify_shape_type(shape) -> str:
    try:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            return "Title"
    except Exception:
        pass
    st = shape.shape_type
    mapping = {
        MSO_SHAPE_TYPE.PICTURE: "Picture",
        MSO_SHAPE_TYPE.TABLE: "Table",
        MSO_SHAPE_TYPE.GROUP: "Group",
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "EmbeddedOLE",
        MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
    }
    return mapping.get(st, "Other")
